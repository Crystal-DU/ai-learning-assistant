from flask import Flask, render_template, request
from pypdf import PdfReader
from pptx import Presentation
import requests
from bs4 import BeautifulSoup

app = Flask(__name__)


def build_prompt(note, mode):
    if mode == "meeting":
        return f"""
# AI Meeting Summary Prompt

You are an AI meeting assistant.

Please analyze the following meeting notes and generate structured meeting minutes.

## Meeting Overview
- Summarize the purpose and main topic of the meeting.

## Key Discussion Points
- List the important discussion points.

## Decisions Made
- List decisions made during the meeting.

## Action Items
- List action items with owner and deadline if available.

## Risks / Issues
- List any risks, blockers, or open issues.

## Next Steps
- List the next steps.

## Meeting Notes

{note}
"""

    return f"""
# AI Learning Summary Prompt

You are an AI learning assistant.

Please analyze the following content and generate a structured summary.

## Topics Learned
- List the main topics.

## Key Takeaways
- Summarize the important points.

## Action Items
- List next actions or things to review.

## Interview Talking Points
- Explain how I can describe this learning experience in a professional setting.

## Content

{note}
"""


@app.route("/", methods=["GET", "POST"])
def index():
    prompt = ""
    raw_text = ""
    extracted_content = ""
    source_name = "Manual Input"
    source_type = "TEXT"
    mode = "prompt"

    if request.method == "POST":
        action = request.form.get("action")
        mode = request.form.get("mode", "prompt")
        raw_text = request.form.get("raw_text", "")
        url = request.form.get("url", "")
        edited_content = request.form.get("extracted_content", "")

        if action == "extract":
            extracted_content = raw_text

            if raw_text:
                source_name = "Manual Input"
                source_type = "TEXT"

            if url:
                source_name = url
                source_type = "URL"

                headers = {"User-Agent": "Mozilla/5.0"}
                response = requests.get(url, headers=headers)
                soup = BeautifulSoup(response.text, "html.parser")

                for tag in soup(["script", "style"]):
                    tag.decompose()

                article = (
                    soup.find("main")
                    or soup.find("article")
                    or soup.find("div", id="mw-content-text")
                )

                if article:
                    extracted_content = article.get_text(separator="\n")
                else:
                    extracted_content = soup.get_text(separator="\n")

                lines = extracted_content.splitlines()
                cleaned_lines = []

                for line in lines:
                    line = line.strip()
                    if len(line) > 40:
                        cleaned_lines.append(line)

                extracted_content = "\n".join(cleaned_lines)

            uploaded_file = request.files.get("file")

            if uploaded_file and uploaded_file.filename:
                source_name = uploaded_file.filename
                filename = uploaded_file.filename.lower()

                if filename.endswith(".pdf"):
                    source_type = "PDF"
                    reader = PdfReader(uploaded_file)
                    extracted_content = ""

                    for page in reader.pages:
                        text = page.extract_text()
                        if text:
                            extracted_content += text + "\n"

                elif filename.endswith(".pptx"):
                    source_type = "PPTX"
                    presentation = Presentation(uploaded_file)
                    extracted_content = ""

                    for slide_number, slide in enumerate(presentation.slides, start=1):
                        extracted_content += f"\n--- Slide {slide_number} ---\n"

                        for shape in slide.shapes:
                            if hasattr(shape, "text"):
                                text = shape.text.strip()
                                if text:
                                    extracted_content += text + "\n"

                elif filename.endswith(".txt"):
                    source_type = "TXT"
                    extracted_content = uploaded_file.read().decode("utf-8")

                elif filename.endswith(".md"):
                    source_type = "MARKDOWN"
                    extracted_content = uploaded_file.read().decode("utf-8")

                else:
                    source_type = "UNKNOWN"
                    extracted_content = uploaded_file.read().decode("utf-8")

        elif action == "generate":
            extracted_content = edited_content
            prompt = build_prompt(extracted_content, mode)

    return render_template(
        "index.html",
        raw_text=raw_text,
        extracted_content=extracted_content,
        prompt=prompt,
        source_name=source_name,
        source_type=source_type,
        mode=mode
    )


if __name__ == "__main__":
    app.run(debug=True)